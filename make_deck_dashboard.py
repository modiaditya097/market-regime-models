"""
Generate 4-slide deck covering dashboard infrastructure.
Matches style of make_deck.py (Arial, black/gray palette, 13.33x7.50").
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
BLACK  = RGBColor(0x00, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
MGRAY  = RGBColor(0xA5, 0xA5, 0xA5)
DGRAY  = RGBColor(0x55, 0x55, 0x55)
RULE   = RGBColor(0x1A, 0x1A, 0x1A)

W, H = 13.33, 7.50

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────
def i(x):
    return Inches(x)

def rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(1, i(l), i(t), i(w), i(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h,
        size=11, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(i(l), i(t), i(w), i(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def hline(slide, l, t, w):
    rect(slide, l, t, w, 0.008, RULE)

def vline(slide, l, t, h):
    s = slide.shapes.add_shape(1, i(l), i(t), i(0.008), i(h))
    s.fill.solid()
    s.fill.fore_color.rgb = RULE
    s.line.fill.background()

def page_num(slide, n):
    rect(slide, 0, H - 0.49, W, 0.49, LGRAY)
    txt(slide, str(n), W - 0.70, H - 0.40, 0.55, 0.32,
        size=11, color=DGRAY, align=PP_ALIGN.RIGHT)

def slide_header(slide, title, subtitle=""):
    txt(slide, title, 0.18, 0.23, 11.58, 0.72, size=27, bold=True, color=BLACK)
    if subtitle:
        txt(slide, subtitle, 0.18, 0.88, 11.58, 0.32, size=13, color=DGRAY)
    hline(slide, 0.18, 1.18, 12.67)

def card(slide, x, y, w, h, title, bullets):
    rect(slide, x, y, w, h, LGRAY)
    txt(slide, title, x + 0.14, y + 0.12, w - 0.26, 0.38,
        size=11, bold=True, color=BLACK)
    hline(slide, x + 0.14, y + 0.52, w - 0.28)
    body = "\n".join(f"  \u2013  {b}" for b in bullets)
    txt(slide, body, x + 0.14, y + 0.60, w - 0.26, h - 0.72,
        size=10, color=DGRAY)


# =============================================================================
# SLIDE 1 — Dashboard Overview
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Dashboard Infrastructure Overview",
             "Posit Shiny for Python  |  Multi-tab reactive UI  |  Config-driven model registration")
page_num(slide, 1)

overview_cards = [
    ("Framework",
     ["Posit Shiny for Python (v0.9+)",
      "Bootstrap dark theme (#1a1a2e)",
      "Entry point: shiny_app/app.py",
      "Run: shiny run shiny_app/app.py --reload"]),
    ("Navigation",
     ["Top-level nav panel with 4 tabs",
      "Tab 1: Cross-model Comparison (fixed)",
      "Tabs 2–4: Per-model tabs (dynamic)",
      "Each tab: sticky sidebar + scrollable main"]),
    ("Model Registry",
     ["Models declared in app_config.yaml",
      "app.py uses importlib to load modules",
      "Add a model = YAML entry + .py module",
      "No changes to app.py required"]),
    ("Key Dependencies",
     ["shiny >= 0.9, pandas >= 2.0",
      "matplotlib >= 3.7, pyyaml >= 6.0",
      "Outputs: CSVs + PNGs per model",
      "No database; file-based state"]),
]

cw, ch = 2.96, 3.10
gap = 0.20
top = 1.32
for idx, (title, bullets) in enumerate(overview_cards):
    x = 0.18 + idx * (cw + gap)
    card(slide, x, top, cw, ch, title, bullets)

hline(slide, 0.18, 4.50, 12.67)
txt(slide,
    "shiny_app/app.py  \u2014  models registered in shiny_app/app_config.yaml  "
    "\u2014  each model module exports model_tab_ui() and model_tab_server()",
    0.18, 4.56, 12.67, 0.30, size=9, color=DGRAY, italic=True)


# =============================================================================
# SLIDE 2 — Tab Structure
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Tab Structure & Per-Model Layout",
             "Four navigation tabs  |  Sidebar controls  |  Reactive TE-target selector")
page_num(slide, 2)

tabs = [
    ("Comparison Tab",
     ["Overlaid cumulative returns (all models)",
      "Side-by-side metrics table",
      "TE-target dropdown (1%–4%)",
      "Placeholder fallback if output missing",
      "Model colors: violet / green / orange"]),
    ("Model 1 — SJM+BL",
     ["Sections: Metrics, Cum. Returns,",
      "  Portfolio Weights, Regime Plots",
      "Run Model button (async, real-time",
      "  progress bar + step counter [1/5])",
      "Live error log (last 20 lines)"]),
    ("Model 2 — Placeholder",
     ["Placeholder card until outputs exist",
      "Optional Run button if run_command",
      "  is set in app_config.yaml",
      "Same module contract as Models 1 & 3",
      "Outputs dir: outputs/model2/"]),
    ("Model 3 — SC-HMM",
     ["9 sections with tabbed sub-views",
      "Performance: Full / In-Sample / OOS",
      "Cum. Returns, Drawdown, Rolling Sharpe",
      "Regime Timeline, Transition Matrix",
      "Stress/Bull tables, Macro Signal cards"]),
]

cw2, ch2 = 2.96, 3.80
for idx, (title, bullets) in enumerate(tabs):
    x = 0.18 + idx * (cw2 + gap)
    card(slide, x, 1.32, cw2, ch2, title, bullets)

hline(slide, 0.18, 5.18, 12.67)
txt(slide,
    "Sidebar (190–200 px) contains: anchor navigation  |  TE-target selector  |  Run button  |  Status indicator",
    0.18, 5.24, 12.67, 0.30, size=9, color=DGRAY, italic=True)


# =============================================================================
# SLIDE 3 — Data Flow
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Data Flow: Pipeline \u2192 Outputs \u2192 Dashboard",
             "Five-step pipeline produces file-based outputs; dashboard renders on demand")
page_num(slide, 3)

# Three-tier flow: raw data | pipeline | outputs | dashboard
tier_w = 2.80
tier_h = 4.40
tier_top = 1.35
tier_gap = 0.22

tiers = [
    ("Raw Data Sources",
     ["Ken French 5-factor returns",
      "Momentum factor (custom)",
      "FRED API: 2Y, 10Y yields",
      "yfinance: VIX, market returns",
      "Parquet cache: outputs/cache/"]),
    ("5-Step Pipeline  (main.py)",
     ["1. data.py — download & cache",
      "2. features.py — 17 feats/factor",
      "     (EWMAs, RSI, MACD, beta)",
      "3. regime.py — SJM fit & infer",
      "4. portfolio.py — BL + SLSQP",
      "5. backtest.py — metrics & plots"]),
    ("Output Contract  (per model)",
     ["results.csv / metrics_*.csv",
      "returns_te{N}.csv (daily returns)",
      "plots/cumulative_returns.png",
      "plots/portfolio_weights.png",
      "plots/regime_*.png, drawdown.png"]),
    ("Dashboard Rendering",
     ["charts.py: PNG \u2192 base64 data URI",
      "charts.py: load_metrics_row()",
      "  filters results.csv by TE",
      "Shiny reactive: re-renders on",
      "  TE selector change (no reload)"]),
]

for idx, (title, bullets) in enumerate(tiers):
    x = 0.18 + idx * (tier_w + tier_gap)
    card(slide, x, tier_top, tier_w, tier_h, title, bullets)
    if idx < 3:
        txt(slide, "\u2192", x + tier_w + 0.04, tier_top + tier_h / 2 - 0.20,
            0.16, 0.40, size=16, bold=True, color=MGRAY, align=PP_ALIGN.CENTER)

hline(slide, 0.18, 5.83, 12.67)
txt(slide,
    "Model 1 config: config.yaml  |  Dashboard config: shiny_app/app_config.yaml  "
    "|  Each model resolves output_dir relative to project root",
    0.18, 5.89, 12.67, 0.30, size=9, color=DGRAY, italic=True)


# =============================================================================
# SLIDE 4 — Extensibility & Module Contract
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Extensibility: Adding a New Model",
             "YAML registration + Python module contract  |  Zero changes to app.py")
page_num(slide, 4)

vline(slide, 6.68, 1.28, 5.10)

# Left column: YAML registration
txt(slide, "Step 1 — Register in app_config.yaml",
    0.18, 1.32, 6.30, 0.36, size=11, bold=True)
hline(slide, 0.18, 1.70, 6.30)

yaml_lines = [
    ("models:", False),
    ("  - id: model4", False),
    ('    name: "My New Model"', False),
    ('    output_dir: "outputs/model4"', False),
    ('    module: "shiny_app.modules.model4"', False),
    ("    te_targets: [0.02, 0.03]", False),
    ('    run_command: ["python", "run_model4.py"]', False),
]
code_bg_w = 6.20
rect(slide, 0.18, 1.78, code_bg_w, 2.18, LGRAY)
for li, (line, bold) in enumerate(yaml_lines):
    txt(slide, line, 0.30, 1.84 + li * 0.30, code_bg_w - 0.20, 0.28,
        size=9.5, bold=bold, color=BLACK)

txt(slide, "app.py auto-discovers and loads the tab via importlib — no edits needed.",
    0.18, 4.04, 6.30, 0.36, size=10, color=DGRAY, italic=True)

# Left — module contract
txt(slide, "Step 2 — Module contract (shiny_app/modules/model4.py)",
    0.18, 4.46, 6.30, 0.36, size=11, bold=True)
hline(slide, 0.18, 4.84, 6.30)
contract = [
    "def model_tab_ui(model_id, cfg):   \u2192 returns ui.TagList",
    "def model_tab_server(model_id, input, output, session, cfg):  \u2192 None",
    "Use cfg['output_dir'] to resolve CSVs and PNGs",
    "Use components/layout.py helpers: section(), placeholder_card()",
]
for li, line in enumerate(contract):
    txt(slide, f"\u2013  {line}", 0.18, 4.92 + li * 0.26, 6.30, 0.24,
        size=9.5, color=DGRAY)

# Right column: output contract table
txt(slide, "Required Output Files",
    6.85, 1.32, 6.20, 0.36, size=11, bold=True)
hline(slide, 6.85, 1.70, 6.20)

output_rows = [
    ("File", "Purpose", True),
    ("metrics_full.csv", "Performance table — full period", False),
    ("metrics_train.csv", "Performance table — in-sample", False),
    ("metrics_test.csv", "Performance table — out-of-sample", False),
    ("returns_te{N}.csv", "Daily returns (date, portfolio, mkt, ew)", False),
    ("plots/cumulative_returns.png", "Cum. returns chart", False),
    ("plots/drawdown.png", "Drawdown chart", False),
    ("plots/portfolio_weights.png", "Weight allocation over time", False),
    ("macro_latest.csv", "Macro signals (indicator, value, signal)", False),
]
col1_w, col2_w = 2.60, 3.42
for ri, (f, purpose, header) in enumerate(output_rows):
    bg = BLACK if header else (LGRAY if ri % 2 == 1 else WHITE)
    fg = WHITE if header else BLACK
    rect(slide, 6.85, 1.76 + ri * 0.36, col1_w + col2_w, 0.36, bg)
    txt(slide, f, 6.90, 1.78 + ri * 0.36, col1_w, 0.30,
        size=9, bold=header, color=fg)
    txt(slide, purpose, 6.90 + col1_w, 1.78 + ri * 0.36, col2_w - 0.10, 0.30,
        size=9, bold=header, color=fg)

txt(slide,
    "Columns in metrics CSVs: Strategy, Cumul.Ret, CAGR, Ann.Vol, Sharpe, Sortino, Max DD, Calmar, Hit Rate",
    6.85, 5.06, 6.20, 0.30, size=8.5, color=DGRAY, italic=True)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "outputs/Dashboard_Infrastructure.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
