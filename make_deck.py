"""
Generate results deck for Dynamic Factor Allocation
Styled after 'Literature Review - Group 1.pptx'
  - 13.33 x 7.50"  white background
  - Arial font, black text
  - 27pt bold title + regular subtitle + thin black horizontal rule
  - Clean 2-column layouts, vertical dividers
  - Light-gray footer strip + page number bottom-right
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ───────────────────────────────────────────────────────────────────
BLACK  = RGBColor(0x00, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)   # footer strip / card bg
MGRAY  = RGBColor(0xA5, 0xA5, 0xA5)   # title-slide accent line
DGRAY  = RGBColor(0x55, 0x55, 0x55)   # secondary text
RULE   = RGBColor(0x1A, 0x1A, 0x1A)   # horizontal/vertical divider lines
ACCENT = RGBColor(0xDF, 0x70, 0x23)   # orange accent (closing slide)

# ── Canvas ────────────────────────────────────────────────────────────────────
W, H = 13.33, 7.50   # inches – matches Literature Review deck

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

PLOTS = "outputs/plots"


# ── Helpers ───────────────────────────────────────────────────────────────────

def i(x):
    return Inches(x)

def rect(slide, l, t, w, h, fill, line=False):
    s = slide.shapes.add_shape(1, i(l), i(t), i(w), i(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = fill
        s.line.width = Emu(9525)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h,
        size=11, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, italic=False, font="Arial"):
    tb = slide.shapes.add_textbox(i(l), i(t), i(w), i(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def hline(slide, l, t, w, thickness=0.008):
    """Thin horizontal rule (solid black)."""
    rect(slide, l, t, w, thickness, RULE)

def vline(slide, l, t, h, thickness=0.008):
    """Thin vertical rule (solid black)."""
    rect(slide, l, t, thickness, h, RULE)

def image(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, i(l), i(t), i(w), i(h))

def page_num(slide, n):
    """Page number bottom-right + light-gray footer strip."""
    rect(slide, 0, H - 0.49, W, 0.49, LGRAY)
    txt(slide, str(n), W - 0.70, H - 0.40, 0.55, 0.32,
        size=11, color=DGRAY, align=PP_ALIGN.RIGHT)

def slide_header(slide, title, subtitle=""):
    """Title + subtitle + horizontal rule — matches content slide layout."""
    txt(slide, title, 0.18, 0.23, 11.58, 0.72,
        size=27, bold=True, color=BLACK)
    if subtitle:
        txt(slide, subtitle, 0.18, 0.88, 11.58, 0.32,
            size=13, color=DGRAY)
    hline(slide, 0.18, 1.18, 12.67)


# =============================================================================
# SLIDE 1 — Title  (Shield layout style)
# =============================================================================
slide = prs.slides.add_slide(BLANK)

# White background (default) + gray accent line on left edge
rect(slide, 0, 0, 0.08, H, MGRAY)

# Footer
rect(slide, 0, H - 0.49, W, 0.49, LGRAY)

# Title block
txt(slide, "QWIM - Market Regimes",
    0.25, 2.36, 9.33, 1.10, size=36, bold=True, color=BLACK)
txt(slide, "Phase 1",
    0.25, 3.82, 7.39, 0.80, size=22, color=DGRAY)
txt(slide, "Aditya Modi\nTanmay Kadam\nIshan Kakodkar",
    0.25, 5.28, 7.41, 1.10, size=16, color=DGRAY)

# Right-side "at a glance" summary (clean, no colored bg)
hline(slide, 9.50, 1.80, 3.50)
txt(slide, "At a Glance", 9.55, 1.90, 3.40, 0.40,
    size=13, bold=True, color=BLACK)
hline(slide, 9.50, 2.32, 3.50)
glance = [
    ("Universe",    "6 factor portfolios (FF5 + Mom)"),
    ("Period",      "Jan 2000 – Jan 2026  (6,511 days)"),
    ("Regime",      "Sparse Jump Model, monthly refit"),
    ("Portfolio",   "Black-Litterman + SLSQP"),
    ("Best Sharpe", "0.57  at TE = 2%  (vs 0.55 EW)"),
]
for idx, (k, v) in enumerate(glance):
    y = 2.40 + idx * 0.56
    txt(slide, k,  9.55, y,       1.20, 0.36, size=11, bold=True)
    txt(slide, v,  10.78, y,      2.40, 0.36, size=11, color=DGRAY)


# =============================================================================
# SLIDE 2 — Methodology Overview
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Methodology",
             "Sparse Jump Model regime detection  +  Black-Litterman portfolio construction")
page_num(slide, 2)

# 5 pipeline step boxes — white cards, bold header, thin rule underneath
step_data = [
    ("1. Feature Engineering",
     "17 features/factor\nEWMA returns (spans 8/21/63d)\nRSI, Stochastic %K, MACD\nLog downside dev, active beta\n4 macro: VIX, 2Y, 10Y–2Y"),
    ("2. Clip + Scale",
     "DataClipperStd (3-sigma)\nStandardScalerPD\nFit on train window only\nNo lookahead in scaling"),
    ("3. SJM Fit",
     "Expanding window (8–12 yr)\nMonthly refit dates\nlambda=50, kappa²=9.5\nn_init_jm=10 restarts"),
    ("4. Online Inference",
     "predict_online() with\n1-day prediction delay\nState 0 = Bull\n(sorted by cumul. return)"),
    ("5. BL Portfolio",
     "5×6 view matrix P\nOmega via Brent's method\nTarget TE: 1/2/3/4%\nSLSQP long-only optimizer"),
]
card_w = 2.42
gap    = 0.10
start  = 0.25
top    = 1.30
card_h = 4.50

for idx, (hdr, body) in enumerate(step_data):
    x = start + idx * (card_w + gap)
    # card bg
    rect(slide, x, top, card_w, card_h, LGRAY)
    # card header — just bold text, thin rule below
    txt(slide, hdr, x + 0.10, top + 0.10, card_w - 0.20, 0.52,
        size=11, bold=True, color=BLACK)
    hline(slide, x + 0.10, top + 0.64, card_w - 0.20)
    # body
    txt(slide, body, x + 0.10, top + 0.72, card_w - 0.20, card_h - 0.90,
        size=10, color=DGRAY)
    # arrow between cards
    if idx < 4:
        txt(slide, ">", x + card_w + 0.01, top + card_h / 2 - 0.20, 0.12, 0.40,
            size=14, bold=True, color=MGRAY, align=PP_ALIGN.CENTER)

# Key design note (bottom strip above footer)
hline(slide, 0.18, 5.82, 12.67)
txt(slide,
    "State 0 = Bull  |  1-day prediction delay  |  Omega calibrated via Brent's method  |  Long-only SLSQP rebalance",
    0.18, 5.88, 12.67, 0.36, size=10, color=DGRAY, italic=True)


# =============================================================================
# SLIDE 3 — Regime Detection Results
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Regime Detection Results",
             "Online inferred bull/bear regimes per factor  |  Test period 2008–2026  |  4,514 trading days")
page_num(slide, 3)

plot_files = [
    ("Value",    f"{PLOTS}/regime_value.png"),
    ("Size",     f"{PLOTS}/regime_size.png"),
    ("Quality",  f"{PLOTS}/regime_quality.png"),
    ("Growth",   f"{PLOTS}/regime_growth.png"),
    ("Momentum", f"{PLOTS}/regime_momentum.png"),
]
regime_stats = [43.5, 44.0, 44.5, 60.1, 60.4]

# 3 across top row, 2 centred in bottom row
positions = [
    (0.15,  1.28, 4.12, 2.10),
    (4.42,  1.28, 4.12, 2.10),
    (8.70,  1.28, 4.12, 2.10),
    (2.10,  3.60, 4.12, 2.10),
    (7.00,  3.60, 4.12, 2.10),
]

for (name, path), (l, t, w, h), bull_pct in zip(plot_files, positions, regime_stats):
    # label above image
    txt(slide, f"{name}  |  {bull_pct:.1f}% Bull / {100-bull_pct:.1f}% Bear",
        l, t - 0.24, w, 0.22, size=10, bold=True, color=BLACK)
    # thin rule above image
    hline(slide, l, t - 0.04, w)
    # image
    image(slide, path, l, t, w, h)


# =============================================================================
# SLIDE 4 — Performance Results
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Empirical Results",
             "Walk-forward backtest 2008–2026  |  After 5 bps/side transaction costs")
page_num(slide, 4)

# ── Left column: KPI summary boxes ──────────────────────────────────────────
vline(slide, 6.52, 1.28, 5.00)

kpis = [
    ("Best Sharpe Ratio",  "0.57",  "Dynamic TE=2%  (vs 0.55 EW benchmark)"),
    ("Ann. Excess Return", "9.70%", "Dynamic TE=2%  (vs 9.51% EW)"),
    ("Max Drawdown",       "-51.9%","Dynamic TE=2%  (vs -53.0% EW)"),
    ("Info Ratio vs EW",   "+0.04", "TE=1% and TE=2% targets"),
]
for idx, (label, big, sub) in enumerate(kpis):
    y = 1.38 + idx * 1.18
    hline(slide, 0.18, y, 6.10)
    txt(slide, label, 0.22, y + 0.06, 3.80, 0.36, size=11, bold=True, color=BLACK)
    txt(slide, big,   0.22, y + 0.40, 3.80, 0.55, size=26, bold=True, color=BLACK)
    txt(slide, sub,   0.22, y + 0.92, 5.80, 0.26, size=10, color=DGRAY)

# ── Right column: performance table ─────────────────────────────────────────
headers  = ["Strategy",       "Ann. Ret", "Sharpe", "Max DD",  "IR vs EW", "Turnover"]
col_w    = [2.65,              0.95,       0.90,     0.95,      0.90,       0.95]
row_data = [
    ("EW Benchmark",    "9.51%", "0.55", "-53.0%", "—",     "—"),
    ("Dynamic  TE=1%",  "9.60%", "0.56", "-53.0%", "+0.04", "2.7x/yr"),
    ("Dynamic  TE=2%",  "9.70%", "0.57", "-51.9%", "+0.04", "5.7x/yr"),
    ("Dynamic  TE=3%",  "9.54%", "0.56", "-52.4%", "-0.02", "9.2x/yr"),
    ("Dynamic  TE=4%",  "9.54%", "0.56", "-52.8%", "-0.02", "11.8x/yr"),
]

def trow(slide, cells, widths, l, t, row_h, bg, fg=BLACK, bold=False, fs=10):
    x = l
    rect(slide, l, t, sum(widths), row_h, bg)
    for cell, cw in zip(cells, widths):
        txt(slide, cell, x + 0.07, t + 0.05, cw - 0.12, row_h - 0.07,
            size=fs, bold=bold, color=fg, align=PP_ALIGN.CENTER)
        x += cw

table_l = 6.65
trow(slide, headers, col_w, table_l, 1.28, 0.42, BLACK, fg=WHITE, bold=True, fs=10)
for idx, row in enumerate(row_data):
    bg   = LGRAY if idx % 2 == 0 else WHITE
    bold = (idx == 2)   # TE=2% best row
    trow(slide, row, col_w, table_l, 1.70 + idx * 0.42, 0.42, bg,
         fg=BLACK, bold=bold, fs=10)

# best-row annotation
txt(slide, "<-- Best risk-adjusted", table_l + sum(col_w) + 0.05, 1.70 + 2 * 0.42 + 0.06,
    1.30, 0.30, size=9, color=DGRAY, italic=True)

hline(slide, 6.65, 3.82, sum(col_w))
txt(slide,
    "All results are out-of-sample  |  Training: expanding window (8 yr min, 12 yr max)  |  SJM refitted monthly",
    6.65, 3.90, sum(col_w), 0.28, size=9, color=DGRAY, italic=True)


# =============================================================================
# SLIDE 5 — Portfolio Analysis
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Portfolio Analysis",
             "Cumulative excess returns and dynamic weight allocation over time")
page_num(slide, 5)

vline(slide, 6.68, 1.28, 5.00)

# Left: cumulative returns
cum_path = f"{PLOTS}/cumulative_returns.png"
txt(slide, "Cumulative Excess Returns", 0.18, 1.28, 6.25, 0.36,
    size=11, bold=True, color=BLACK)
hline(slide, 0.18, 1.62, 6.25)
image(slide, cum_path, 0.18, 1.68, 6.28, 3.60)
txt(slide, "Solid = Dynamic strategies  |  Dashed = EW benchmark  |  Dotted = Market",
    0.18, 5.34, 6.28, 0.28, size=9, color=DGRAY, italic=True)

# Right: portfolio weights
wgt_path = f"{PLOTS}/portfolio_weights.png"
txt(slide, "Dynamic Portfolio Weights  (TE=3%)", 6.80, 1.28, 6.20, 0.36,
    size=11, bold=True, color=BLACK)
hline(slide, 6.80, 1.62, 6.20)
image(slide, wgt_path, 6.80, 1.68, 6.20, 3.60)
txt(slide, "Regime-driven tilts shift weight across factors over time",
    6.80, 5.34, 6.20, 0.28, size=9, color=DGRAY, italic=True)


# =============================================================================
# SLIDE 6 — Key Takeaways
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Key Takeaways",
             "Summary of findings from Phase 1 replication")
page_num(slide, 6)

takeaways = [
    ("Regime detection is reliable",
     "SJM identifies persistent bull/bear states per factor. All 5 factors show clear regime "
     "structure. Growth and Momentum spend 60%+ in bull regimes; Value/Size/Quality are more "
     "bear-heavy over the 2008–2026 test period."),
    ("Low-TE tilts deliver value",
     "TE=2% achieves best Sharpe (0.57 vs 0.55 EW) and positive IR (+0.04) with manageable "
     "turnover (5.7x/year). The BL framework correctly translates regime signals into modest, "
     "cost-aware tilts."),
    ("Aggressive tilts hurt after costs",
     "TE=3–4% targets increase turnover to 9–12x/year without commensurate return improvement. "
     "Transaction costs erode the alpha from stronger tilts. Sweet spot is TE = 1–2%."),
    ("Signal quality is the bottleneck",
     "The regime signal is informative but not strong enough to support large active bets. "
     "Improvements should focus on richer features, better hyperparameter tuning, and "
     "expanding the factor universe."),
]

card_w2 = 6.10
card_h2 = 2.18
for idx, (hdr, body) in enumerate(takeaways):
    col = idx % 2
    row = idx // 2
    x = 0.18 + col * (card_w2 + 0.60)
    y = 1.30 + row * (card_h2 + 0.18)
    rect(slide, x, y, card_w2, card_h2, LGRAY)
    txt(slide, hdr,  x + 0.16, y + 0.12, card_w2 - 0.30, 0.38,
        size=12, bold=True, color=BLACK)
    hline(slide, x + 0.16, y + 0.52, card_w2 - 0.30)
    txt(slide, body, x + 0.16, y + 0.60, card_w2 - 0.30, card_h2 - 0.70,
        size=10, color=DGRAY)


# =============================================================================
# SLIDE 7 — Next Steps
# =============================================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Phase 2 Roadmap",
             "Contributions and enhancements beyond Phase 1 replication")
page_num(slide, 7)

next_steps = [
    ("01 — Hyperparameter Optimization",
     "Grid/Bayesian search over SJM lambda (jump penalty) and kappa² (sparsity)\n"
     "Cross-validate TE targets on a held-out validation set (avoid overfitting to 2%)\n"
     "Per-factor tuning — value vs momentum may need different regime granularity\n"
     "Sensitivity analysis: stability of results to parameter perturbations"),
    ("02 — Expanded Factor Universe",
     "Add Low Volatility (BAB/IVOL from AQR or constructed from daily returns)\n"
     "Quality decomposition: profitability vs investment vs accruals as separate factors\n"
     "International factors: EM and DM versions of FF5 (available from Ken French)\n"
     "Sector momentum and industry-level factors for finer regime granularity"),
    ("03 — Model Enhancements",
     "Compare SJM against HMM-MLE and Gaussian Mixture baselines\n"
     "Explore 3-state models (Bull / Neutral / Bear) for richer regime structure\n"
     "Ensemble regimes: combine per-factor signals into market-wide indicator\n"
     "Online/incremental SJM updates to reduce monthly refit latency"),
    ("04 — Portfolio & Infrastructure",
     "Add turnover constraints directly into SLSQP (hard limit on rebalance size)\n"
     "Risk-parity and factor-risk-parity as alternative priors to equal-weight\n"
     "Streamlit dashboard: live regime monitor + parameter explorer\n"
     "Stress-test on sub-periods: GFC (2008), COVID (2020), rate hikes (2022)"),
]

for idx, (hdr, body) in enumerate(next_steps):
    col = idx % 2
    row = idx // 2
    x = 0.18 + col * (card_w2 + 0.60)
    y = 1.30 + row * (card_h2 + 0.18)
    rect(slide, x, y, card_w2, card_h2, LGRAY)
    txt(slide, hdr, x + 0.16, y + 0.12, card_w2 - 0.30, 0.38,
        size=12, bold=True, color=BLACK)
    hline(slide, x + 0.16, y + 0.52, card_w2 - 0.30)
    for bi, line in enumerate(body.strip().split("\n")):
        txt(slide, line.strip(), x + 0.16, y + 0.62 + bi * 0.36,
            card_w2 - 0.30, 0.34, size=10, color=DGRAY)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "outputs/Dynamic_Factor_Allocation_Results_v2.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
